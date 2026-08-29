"""ParseDoc Text Extraction helpers"""

import re
from typing import Dict, List, Tuple


# Common OCR / typographical errors seen in regulatory Gazette PDFs.
# Conservative: only exact, unambiguous misspellings are corrected.
_PDF_TYPO_FIXES = {
    "regulatiions": "regulations",
    "hydorxide": "hydroxide",
    "Boudouins": "Boudouin's",
    "rancidityand": "rancidity and",
    "basisBasis": "basis Basis",
    "Adminidive": "Andaman",
    "Orisssa": "Orissa",
    "Kerela": "Kerala",
    "Lakshwadeep": "Lakshadweep",
    "Meghalya": "Meghalaya",
    "AssamBihar": "Assam, Bihar",
}


def _page_text_words(page) -> str:
    """Reconstruct a page's text from word bounding boxes.

    PyMuPDF's plain ``get_text("text")`` drops spaces in justified/space-less
    PDFs (e.g. "AssamBihar"). Using ``get_text("words")`` and joining by
    x/y position restores word spacing far more reliably.
    """
    words = page.get_text("words")  # x0, y0, x1, y1, word, block, line, wno
    if not words:
        return page.get_text("text") or ""
    words.sort(key=lambda w: (round(w[1]), w[0]))
    lines: List[List] = []
    current: List = []
    current_y: int = None
    for w in words:
        y = round(w[1])
        if current_y is None:
            current_y = y
        if abs(y - current_y) <= 2:
            current.append(w)
        else:
            lines.append(current)
            current_y = y
            current = [w]
    if current:
        lines.append(current)
    return "\n".join(" ".join(x[4] for x in ln) for ln in lines)


def _fix_pdf_text_artifacts(text: str) -> str:
    """Repair common PyMuPDF extraction artifacts for regulatory PDFs.

    - The degree sign is often extracted as a stray "0" (63°C -> 630 C).
    - Concatenated legal phrases ("Notmorethan" -> "Not more than").
    - Stray LaTeX-ish math markup from structured PDFs.
    - A small set of unambiguous OCR typos.
    """
    # Strip math delimiters and repair degree superscripts like $115^{0}C$.
    # Also swallows an optional trailing " C" that some PDFs emit after the close.
    text = re.sub(
        r"\$\s*(\d+(?:\.\d+)?)\s*\^\s*\{\s*0\s*\}\s*([Cc])\s*\$(?:\s*[Cc])?", r"\1°C", text
    )
    text = text.replace("\\times", "×")
    text = text.replace("$", " ")
    # 63°C -> "630 C" / "71.5°C" -> "71.50C" / "10°C" -> "100 C" etc.
    # Only fire when there is no pre-existing degree sign (so "130°C" is left alone).
    text = re.sub(r"(\d)0\s*C\b", r"\1°C", text, flags=re.IGNORECASE)
    text = text.replace("Notmorethan", "Not more than")
    text = text.replace("Notlessthan", "Not less than")
    # Spacing between a number and a unit word.
    text = re.sub(
        r"(\d)\s*(percent|per\s*cent|mg|µg|mcg|g|kg|ml|l|ppm)\b",
        r"\1 \2",
        text,
        flags=re.IGNORECASE,
    )
    # Unambiguous OCR typos.
    for bad, good in _PDF_TYPO_FIXES.items():
        text = text.replace(bad, good)
    # Collapse the double spaces introduced above.
    text = re.sub(r" {2,}", " ", text)
    return text


def _extract_page_tables(page) -> Tuple[List[List[List[str]]], List[Tuple[float, float]]]:
    """Extract column-aligned tables from a page using word positions.

    Many regulatory PDFs lay tables out as whitespace-separated columns with no
    vector rules, so PyMuPDF's ``find_tables`` cannot see them. We detect runs
    of consecutive lines that each contain >=2 aligned columns and reconstruct
    rows. Returns (tables, skipped_y_ranges) where ``skipped_y_ranges`` are the
    (y0, y1) bands covered by tables so paragraph text can exclude them.
    """
    words = page.get_text("words")  # x0, y0, x1, y1, word, block, line, wno
    if not words:
        return [], []
    # Cluster word start-x positions into column lanes.
    x0s = sorted({round(w[0], 1) for w in words})
    if not x0s:
        return [], []
    col_gap = 22.0
    clusters: List[List[float]] = [[x0s[0]]]
    for x in x0s[1:]:
        if x - clusters[-1][-1] > col_gap:
            clusters.append([x])
        else:
            clusters[-1].append(x)
    col_centers = [sum(c) / len(c) for c in clusters]

    # Group words into lines by y.
    words.sort(key=lambda w: (round(w[1]), w[0]))
    lines: List[List] = []
    cur: List = []
    cur_y: int = None
    for w in words:
        y = round(w[1])
        if cur_y is None:
            cur_y = y
        if abs(y - cur_y) <= 3:
            cur.append(w)
        else:
            lines.append(cur)
            cur_y = y
            cur = [w]
    if cur:
        lines.append(cur)

    def assign_cols(line_words):
        row = {}
        for w in line_words:
            cx = (w[0] + w[2]) / 2.0
            best = min(range(len(col_centers)), key=lambda i: abs(col_centers[i] - cx))
            if abs(col_centers[best] - cx) <= col_gap:
                row.setdefault(best, []).append(w[4])
        return row

    tables: List[List[List[str]]] = []
    skipped: List[Tuple[float, float]] = []
    run: List[List] = []
    run_y: Tuple[float, float] = (1e9, -1e9)
    for ln in lines:
        row = assign_cols(ln)
        ncols = len(row)
        if ncols >= 2:
            ordered = [(" ".join(row.get(i, []))).strip() for i in range(len(col_centers))]
            run.append(ordered)
            run_y = (min(run_y[0], min(w[1] for w in ln)), max(run_y[1], max(w[3] for w in ln)))
        else:
            if len(run) >= 4:
                tables.append(run)
                skipped.append(run_y)
            run = []
            run_y = (1e9, -1e9)
    if len(run) >= 4:
        tables.append(run)
        skipped.append(run_y)
    return tables, skipped


def _paragraph_text_excluding(page, skipped: List[Tuple[float, float]]) -> str:
    """Reconstruct paragraph text, skipping y-bands covered by tables."""
    words = page.get_text("words")
    if not words:
        return page.get_text("text") or ""
    lines: List[List] = []
    cur: List = []
    cur_y: int = None
    for w in sorted(words, key=lambda w: (round(w[1]), w[0])):
        y = round(w[1])
        in_table = any(y0 - 3 <= w[1] <= y1 + 3 for (y0, y1) in skipped)
        if in_table:
            continue
        if cur_y is None:
            cur_y = y
        if abs(y - cur_y) <= 2:
            cur.append(w)
        else:
            lines.append(cur)
            cur_y = y
            cur = [w]
    if cur:
        lines.append(cur)
    return "\n".join(" ".join(x[4] for x in ln) for ln in lines)


def extract_pdf_text(pdf_path: str) -> Dict[str, object]:
    """Extract text and page layout from a PDF using PyMuPDF.

    Returns a dict with ``text``, ``layout``, and ``page_count`` keys.
    """
    try:
        import pymupdf as fitz
    except ImportError:
        import fitz  # PyMuPDF

    text_parts: List[str] = []
    layout_parts: List[str] = []
    all_tables: List[List[List[str]]] = []
    with fitz.open(pdf_path) as doc:
        page_count = doc.page_count
        for i, page in enumerate(doc):
            tables, skipped = _extract_page_tables(page)
            para = _fix_pdf_text_artifacts(_paragraph_text_excluding(page, skipped))
            text_parts.append(para)
            layout_parts.append(f"Page {i + 1}")
            for tbl in tables:
                fixed_rows = [[_fix_pdf_text_artifacts(c) for c in row] for row in tbl]
                all_tables.append(fixed_rows)
    return {
        "text": "\n".join(text_parts).strip(),
        "layout": "\n".join(layout_parts),
        "page_count": page_count,
        "tables": all_tables,
    }


def extract_docx_text(docx_path: str) -> Dict[str, object]:
    """Extract paragraphs and tables from a DOCX using python-docx."""
    from docx import Document

    doc = Document(docx_path)
    text_parts: List[str] = [p.text for p in doc.paragraphs]
    tables: List[List[List[str]]] = []
    for table in doc.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        tables.append(rows)
    return {"text": "\n".join(text_parts).strip(), "tables": tables}


def _pptx_looks_like_list(shape) -> bool:
    """Heuristic: does this text frame contain a bullet/numbered list?"""
    paragraphs = shape.text_frame.paragraphs
    total = 0
    bullets = 0
    for p in paragraphs:
        t = (p.text or "").strip()
        if not t:
            continue
        total += 1
        if (p.level and p.level > 0) or t[0] in "•*-–◦·":
            bullets += 1
    return total > 0 and bullets >= max(1, total // 2)


def _pptx_clean_item(line: str) -> str:
    return line.strip("•*-–◦· \t")


def extract_pptx_text(pptx_path: str) -> Dict[str, object]:
    """Extract structured slide content (titles, lists, tables) from a PPTX."""
    from pptx import Presentation

    prs = Presentation(pptx_path)
    blocks: List[Dict] = []
    text_parts: List[str] = []
    tables: List[List[List[str]]] = []
    first_title: str = ""

    for i, slide in enumerate(prs.slides, 1):
        slide_title = None
        body = []
        for shape in slide.shapes:
            if shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                tables.append(rows)
                blocks.append(
                    {
                        "type": "table",
                        "headers": rows[0] if rows else [],
                        "rows": rows[1:] if len(rows) > 1 else [],
                    }
                )
                continue
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text.strip()
            if not txt:
                continue
            try:
                is_title = bool(
                    shape.is_placeholder and shape.placeholder_format.type == 1  # TITLE
                )
            except Exception:
                is_title = False
            if is_title and slide_title is None:
                slide_title = txt
                continue
            body.append((shape, txt))

        # Fall back to first text frame as the slide title.
        if slide_title is None and body:
            slide_title = body[0][1]
            body = body[1:]

        title = (slide_title or f"Slide {i}").strip()
        # Drop body items that merely repeat the slide title.
        body = [(s, t) for (s, t) in body if t.strip() != title]
        if not first_title:
            first_title = title
        blocks.append({"type": "heading", "level": 1, "text": title})
        text_parts.append(f"# {title}")

        for shape, txt in body:
            if _pptx_looks_like_list(shape):
                items = [_pptx_clean_item(ln) for ln in txt.splitlines() if ln.strip()]
                if items:
                    blocks.append({"type": "list", "ordered": False, "items": items})
                    text_parts.append("")
                    text_parts.extend(f"- {it}" for it in items)
            else:
                blocks.append({"type": "paragraph", "text": txt})
                text_parts.append(txt)
        text_parts.append("")  # slide separator

    return {
        "text": "\n".join(text_parts).strip(),
        "tables": tables,
        "blocks": blocks,
        "title": first_title,
    }


def extract_html_text(html_path: str) -> Dict[str, object]:
    """Extract headings, paragraphs, and tables from HTML using BeautifulSoup."""
    from bs4 import BeautifulSoup

    with open(html_path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "html.parser")

    text_parts: List[str] = []
    for tag in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li"]):
        text_parts.append(tag.get_text(strip=True))
    tables: List[List[List[str]]] = []
    for table in soup.find_all("table"):
        rows = []
        for tr in table.find_all("tr"):
            cells = [c.get_text(strip=True) for c in tr.find_all(["td", "th"])]
            if cells:
                rows.append(cells)
        if rows:
            tables.append(rows)
    return {"text": "\n".join(text_parts).strip(), "tables": tables}


def extract_plain_text(text_path: str) -> Dict[str, object]:
    """Read a plain text or markdown file."""
    with open(text_path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    return {"text": content.strip(), "tables": []}
